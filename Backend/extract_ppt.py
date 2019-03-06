from pptx import Presentation
import os

cur_directory = os.path.dirname(os.path.abspath( __file__ ))
files = [x for x in os.listdir(cur_directory) if x.endswith(".pptx")]

raw_text = []
def ppt_to_text():
    for eachfile in files:
        prs = Presentation(eachfile)
        print(eachfile)
        print("----------------------")
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    raw_text.append(shape.text)
    return str(raw_text)
    